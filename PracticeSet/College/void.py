from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Dark Modern Theme Colors
bg_color = RGBColor(16, 24, 32)   # #101820
accent_cyan = RGBColor(0, 229, 255)
text_color = RGBColor(181, 181, 181)

# Helper to add slide
def add_slide(title, content):
    slide_layout = prs.slide_layouts[6]  # blank slide
    slide = prs.slides.add_slide(slide_layout)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = accent_cyan
    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = content
    p2.font.size = Pt(20)
    p2.font.color.rgb = text_color

# --- Cover Slide ---
add_slide("✨ LINE DRAWING & DDA ALGORITHM",
          "Incremental Rasterization in Computer Graphics\nBy: Tushar Kharade\nRoll No.: IC-2K23-96")

# --- Slide 2 Example ---
add_slide("📌 Introduction to Line Drawing & DDA Algorithm",
          "• Computer screens consist of tiny pixels.\n• A straight line must be represented using selected pixels.\n• Line drawing algorithms convert mathematical lines into pixel coordinates.\n\n• DDA = Digital Differential Analyzer\n• Incremental line drawing algorithm\n• Computes intermediate points using simple addition.\n\nApplications:\n• Computer graphics\n• CAD tools\n• Games\n• Simulations\n• Wireframe models")

# You can continue adding slides here the same way...

# Save presentation
prs.save("DDA_Presentation.pptx")
print("Presentation created successfully!")
